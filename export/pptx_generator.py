"""
Gerador de apresentação PPTX customizada (16 slides).
Cria a apresentação programaticamente com python-pptx.
"""
import io
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from models.inputs import (
    ClienteBasicInfo,
    ProcessoAtual,
    DoresSelecionadas,
    InvestimentoAutomacao,
)
from models.results import ResultadosFinanceiros, MetasReducao

# Paleta de cores
AZUL_ESCURO = RGBColor(0x1F, 0x4E, 0x79)
AZUL_MEDIO = RGBColor(0x2E, 0x75, 0xB6)
AZUL_CLARO = RGBColor(0x9D, 0xC3, 0xE6)
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
CINZA_ESCURO = RGBColor(0x33, 0x33, 0x33)
CINZA_MEDIO = RGBColor(0x77, 0x77, 0x77)
VERDE = RGBColor(0x00, 0xB0, 0x50)
VERMELHO = RGBColor(0xC0, 0x00, 0x00)
LARANJA = RGBColor(0xED, 0x7D, 0x31)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


class PPTXGenerator:
    """Gerador de apresentação PPTX customizada."""

    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_WIDTH
        self.prs.slide_height = SLIDE_HEIGHT

    def gerar(
        self,
        cliente: ClienteBasicInfo,
        processo: ProcessoAtual,
        dores: DoresSelecionadas,
        resultados: ResultadosFinanceiros,
        metas: MetasReducao,
        investimento: InvestimentoAutomacao,
    ) -> io.BytesIO:
        """Gera PPTX completo e retorna como BytesIO."""
        self._slide_01_capa(cliente)
        self._slide_02_agenda()
        self._slide_03_contexto(cliente)
        self._slide_04_processo_atual(processo)
        self._slide_05_dados_operacionais(processo, cliente)
        self._slide_06_analise_estrategica(dores)
        self._slide_07_cenario_critico(resultados)
        self._slide_08_custos_operacionais(resultados)  # Dor 1
        self._slide_09_custos_qualidade(resultados)  # Dor 2
        self._slide_10_custos_seguranca(resultados)  # Dor 3
        self._slide_11_custos_produtividade(resultados)  # Dor 4
        self._slide_12_custos_ocultos(resultados)  # Dor 5
        self._slide_13_consolidacao(resultados)
        self._slide_13_escopo_tecnico()
        self._slide_14_investimento(investimento)
        self._slide_15_viabilidade(resultados, investimento)
        self._slide_16_proximas_etapas()

        buffer = io.BytesIO()
        self.prs.save(buffer)
        buffer.seek(0)
        return buffer

    # =========================================================================
    # Helpers
    # =========================================================================

    def _add_slide(self):
        """Adiciona slide em branco."""
        layout = self.prs.slide_layouts[6]  # blank
        return self.prs.slides.add_slide(layout)

    def _add_textbox(self, slide, left, top, width, height, text,
                     font_size=18, bold=False, color=CINZA_ESCURO,
                     alignment=PP_ALIGN.LEFT, font_name="Calibri"):
        """Adiciona caixa de texto ao slide."""
        txbox = slide.shapes.add_textbox(left, top, width, height)
        tf = txbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = font_name
        p.alignment = alignment
        return txbox

    def _add_title_bar(self, slide, title_text):
        """Adiciona barra de título azul no topo do slide."""
        # Fundo azul
        shape = slide.shapes.add_shape(
            1, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.2),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = AZUL_ESCURO
        shape.line.fill.background()

        # Texto do título
        self._add_textbox(
            slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
            title_text, font_size=28, bold=True, color=BRANCO,
        )

    def _add_subtitle(self, slide, text, top=Inches(1.4)):
        """Adiciona subtítulo abaixo da barra."""
        self._add_textbox(
            slide, Inches(0.5), top, Inches(12), Inches(0.5),
            text, font_size=16, color=CINZA_MEDIO,
        )

    def _add_metric_box(self, slide, left, top, width, height,
                        label, value, color=AZUL_ESCURO):
        """Adiciona caixa de métrica com label e valor."""
        # Borda
        shape = slide.shapes.add_shape(1, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = BRANCO
        shape.line.color.rgb = color
        shape.line.width = Pt(2)

        # Label
        self._add_textbox(
            slide, left + Inches(0.15), top + Inches(0.1),
            width - Inches(0.3), Inches(0.4),
            label, font_size=11, color=CINZA_MEDIO,
            alignment=PP_ALIGN.CENTER,
        )
        # Valor
        self._add_textbox(
            slide, left + Inches(0.15), top + Inches(0.45),
            width - Inches(0.3), Inches(0.6),
            value, font_size=20, bold=True, color=color,
            alignment=PP_ALIGN.CENTER,
        )

    def _add_table(self, slide, left, top, width, height, rows, cols,
                   data, col_widths=None):
        """Adiciona tabela ao slide."""
        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = table_shape.table

        if col_widths:
            for i, w in enumerate(col_widths):
                table.columns[i].width = w

        for r in range(rows):
            for c in range(cols):
                cell = table.cell(r, c)
                cell.text = str(data[r][c])
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(11)
                    paragraph.font.name = "Calibri"
                    if r == 0:
                        paragraph.font.bold = True
                        paragraph.font.color.rgb = BRANCO
                        paragraph.alignment = PP_ALIGN.CENTER
                    else:
                        paragraph.font.color.rgb = CINZA_ESCURO
                        if c > 0:
                            paragraph.alignment = PP_ALIGN.RIGHT

                # Header row styling
                if r == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = AZUL_ESCURO
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = (
                        RGBColor(0xF2, 0xF2, 0xF2) if r % 2 == 0
                        else BRANCO
                    )

        return table

    def _fmt(self, valor: float) -> str:
        """Formata valor monetário."""
        return f"R$ {valor:,.2f}".replace(",", "X").replace(".", ",").replace("X", ".")

    def _fmt_pct(self, valor: float) -> str:
        """Formata percentual."""
        return f"{valor:.1f}%"

    # =========================================================================
    # Slides
    # =========================================================================

    def _slide_01_capa(self, cliente: ClienteBasicInfo):
        slide = self._add_slide()

        # Fundo azul escuro
        bg = slide.shapes.add_shape(
            1, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT,
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = AZUL_ESCURO
        bg.line.fill.background()

        # Título principal
        self._add_textbox(
            slide, Inches(1), Inches(2), Inches(11), Inches(1),
            "Análise do Custo da Inação",
            font_size=40, bold=True, color=BRANCO, alignment=PP_ALIGN.CENTER,
        )

        # Subtítulo
        self._add_textbox(
            slide, Inches(1), Inches(3.2), Inches(11), Inches(0.6),
            "Quanto custa NÃO automatizar?",
            font_size=24, color=AZUL_CLARO, alignment=PP_ALIGN.CENTER,
        )

        # Linha separadora
        line = slide.shapes.add_shape(
            1, Inches(4), Inches(4.1), Inches(5), Inches(0.03),
        )
        line.fill.solid()
        line.fill.fore_color.rgb = AZUL_CLARO
        line.line.fill.background()

        # Info cliente
        self._add_textbox(
            slide, Inches(1), Inches(4.5), Inches(11), Inches(0.5),
            f"Cliente: {cliente.nome_cliente}",
            font_size=20, color=BRANCO, alignment=PP_ALIGN.CENTER,
        )
        self._add_textbox(
            slide, Inches(1), Inches(5.1), Inches(11), Inches(0.5),
            f"Projeto: {cliente.nome_projeto}",
            font_size=18, color=AZUL_CLARO, alignment=PP_ALIGN.CENTER,
        )

        # Data
        data_str = datetime.now().strftime("%d/%m/%Y")
        self._add_textbox(
            slide, Inches(1), Inches(6.2), Inches(11), Inches(0.4),
            data_str, font_size=14, color=AZUL_CLARO, alignment=PP_ALIGN.CENTER,
        )

    def _slide_02_agenda(self):
        slide = self._add_slide()
        self._add_title_bar(slide, "Agenda")

        itens = [
            "1. Contexto e Objetivo",
            "2. Processo Atual",
            "3. Análise Estratégica de Dores",
            "4. Quantificação de Custos",
            "5. Consolidação Financeira",
            "6. Escopo Técnico da Solução",
            "7. Investimento e Viabilidade",
            "8. Próximas Etapas",
        ]

        for i, item in enumerate(itens):
            self._add_textbox(
                slide,
                Inches(1.5), Inches(1.8 + i * 0.6),
                Inches(10), Inches(0.5),
                item, font_size=20, color=AZUL_ESCURO if i % 2 == 0 else CINZA_ESCURO,
            )

    def _slide_03_contexto(self, cliente: ClienteBasicInfo):
        slide = self._add_slide()
        self._add_title_bar(slide, "Contexto e Objetivo")

        self._add_textbox(
            slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(1),
            f"O presente estudo tem como objetivo quantificar o impacto financeiro "
            f"do Custo da Inação no processo de {cliente.nome_projeto} "
            f"e apoiar a decisão de investimento em automação industrial.",
            font_size=16, color=CINZA_ESCURO,
        )

        self._add_textbox(
            slide, Inches(0.8), Inches(3.2), Inches(5), Inches(0.4),
            "Área de atuação ARV:", font_size=14, bold=True, color=AZUL_ESCURO,
        )
        self._add_textbox(
            slide, Inches(6), Inches(3.2), Inches(5), Inches(0.4),
            cliente.area_atuacao, font_size=14, color=CINZA_ESCURO,
        )

        self._add_textbox(
            slide, Inches(0.8), Inches(3.8), Inches(5), Inches(0.4),
            "Porte da empresa:", font_size=14, bold=True, color=AZUL_ESCURO,
        )
        self._add_textbox(
            slide, Inches(6), Inches(3.8), Inches(5), Inches(0.4),
            cliente.porte_empresa, font_size=14, color=CINZA_ESCURO,
        )

        self._add_textbox(
            slide, Inches(0.8), Inches(4.4), Inches(5), Inches(0.4),
            "Fator de encargos:", font_size=14, bold=True, color=AZUL_ESCURO,
        )
        self._add_textbox(
            slide, Inches(6), Inches(4.4), Inches(5), Inches(0.4),
            f"{cliente.fator_encargos:.2f}x", font_size=14, color=CINZA_ESCURO,
        )

        self._add_textbox(
            slide, Inches(0.8), Inches(5.1), Inches(11.5), Inches(1.3),
            "Metodologia:\n"
            "• Levantamento de dados operacionais junto ao cliente\n"
            "• Quantificação por 5 Dores e 18 fórmulas (F01–F18)\n"
            "• Definição de metas de redução com automação\n"
            "• Cálculo de ROI e payback do investimento",
            font_size=14, color=CINZA_ESCURO,
        )

    def _slide_04_processo_atual(self, processo: ProcessoAtual):
        slide = self._add_slide()
        self._add_title_bar(slide, "Processo Atual — Parâmetros de Produção")

        prod_txt = (
            f"{processo.producao_mensal:,.0f} peças/mês" if processo.producao_mensal is not None
            else f"{processo.cadencia_producao} peças/min"
        )

        dados = [
            ("Produção informada", prod_txt),
            ("Horas por Turno", f"{processo.horas_por_turno}h"),
            ("Turnos por Dia", str(processo.turnos_por_dia)),
            ("Dias de Operação por Ano", str(processo.dias_operacao_ano)),
            ("Operadores no Processo (por turno)", str(processo.pessoas_processo_turno)),
            ("Inspetores (por turno)", str(processo.pessoas_inspecao_turno)),
            ("Salário Médio Operador", self._fmt(processo.salario_medio_operador)),
            ("Custo Matéria-Prima/Peça", self._fmt(processo.custo_materia_prima_peca)),
            ("Faturamento Mensal da Linha", self._fmt(processo.faturamento_mensal_linha or 0.0)),
        ]

        table_data = [["Parâmetro", "Valor"]] + [[d[0], d[1]] for d in dados]
        self._add_table(
            slide, Inches(2), Inches(1.8), Inches(9), Inches(4),
            len(table_data), 2, table_data,
            col_widths=[Inches(5.5), Inches(3.5)],
        )

    def _slide_05_dados_operacionais(self, processo: ProcessoAtual, cliente: ClienteBasicInfo):
        slide = self._add_slide()
        self._add_title_bar(slide, "Bases de Cálculo (V2.0)")

        if processo.producao_mensal is not None and processo.producao_mensal > 0:
            prod_anual = processo.producao_mensal * 12
        else:
            cad = processo.cadencia_producao or 0.0
            prod_anual = cad * 60 * processo.horas_por_turno * processo.turnos_por_dia * processo.dias_operacao_ano

        horas_anuais = processo.horas_por_turno * processo.turnos_por_dia * processo.dias_operacao_ano
        pessoas_total = processo.pessoas_processo_turno * processo.turnos_por_dia
        custo_hora_operador = (processo.salario_medio_operador * cliente.fator_encargos) / 176
        custo_hora_parada = (processo.faturamento_mensal_linha or 0.0) / 176

        box_w = Inches(3.5)
        box_h = Inches(1.2)
        gap = Inches(0.5)
        start_x = Inches(0.8)

        self._add_metric_box(
            slide, start_x, Inches(2), box_w, box_h,
            "Produção Anual", f"{prod_anual:,.0f} peças",
        )
        self._add_metric_box(
            slide, start_x + box_w + gap, Inches(2), box_w, box_h,
            "Horas Anuais de Operação", f"{horas_anuais:,.0f} horas",
        )
        self._add_metric_box(
            slide, start_x + 2 * (box_w + gap), Inches(2), box_w, box_h,
            "Pessoas Expostas (Total)", f"{pessoas_total} pessoas",
        )

        self._add_metric_box(
            slide, start_x, Inches(3.8), box_w, box_h,
            "Custo Hora Operador (c/ encargos)", self._fmt(custo_hora_operador),
        )
        self._add_metric_box(
            slide, start_x + box_w + gap, Inches(3.8), box_w, box_h,
            "Custo Hora Parada", self._fmt(custo_hora_parada),
        )
        self._add_metric_box(
            slide, start_x + 2 * (box_w + gap), Inches(3.8), box_w, box_h,
            "Turnos x Dias/Ano",
            f"{processo.turnos_por_dia} turnos x {processo.dias_operacao_ano} dias",
        )

    def _slide_06_analise_estrategica(self, dores: DoresSelecionadas):
        slide = self._add_slide()
        self._add_title_bar(slide, "Análise Estratégica — Dores Identificadas")
        self._add_subtitle(slide, "Custos mapeados no cenário atual do cliente")

        def _lines(itens):
            out = []
            for nome, ativo in itens:
                marcador = "●" if ativo else "○"
                out.append(f"{marcador} {nome}")
            return "\n".join(out)

        dor1 = [
            ("F01: Mão de Obra Direta", dores.f01_mao_de_obra_direta),
            ("F02: Horas Extras", dores.f02_horas_extras),
            ("F03: Curva de Aprendizagem", dores.f03_curva_aprendizagem),
            ("F04: Turnover", dores.f04_turnover),
        ]
        dor2 = [
            ("F05: Refugo e Retrabalho", dores.f05_refugo_retrabalho),
            ("F06: Inspeção Manual", dores.f06_inspecao_manual),
            ("F07: Escapes de Qualidade", dores.f07_escapes_qualidade),
        ]
        dor3 = [
            ("F08: Custo de Oportunidade", dores.f08_custo_oportunidade),
            ("F09: Ociosidade Silenciosa", dores.f09_ociosidade_silenciosa),
            ("F10: Paradas de Linha", dores.f10_paradas_linha),
            ("F11: Setup/Changeover", dores.f11_setup_changeover),
        ]
        dor4 = [
            ("F12: Riscos/Acidentes", dores.f12_riscos_acidentes),
            ("F13: Frota Empilhadeiras", dores.f13_frota_empilhadeiras),
        ]
        dor5 = [
            ("F14: Supervisão", dores.f14_supervisao),
            ("F15: Compliance/EPIs", dores.f15_compliance_epis),
            ("F16: Energia/Utilidades", dores.f16_energia_utilidades),
            ("F17: Espaço Físico", dores.f17_espaco_fisico),
            ("F18: Gestão de Dados", dores.f18_gestao_dados),
        ]

        self._add_textbox(slide, Inches(0.8), Inches(2.1), Inches(5.8), Inches(2.4), "Dor 1: Mão de Obra", font_size=14, bold=True, color=AZUL_ESCURO)
        self._add_textbox(slide, Inches(0.8), Inches(2.5), Inches(5.8), Inches(1.8), _lines(dor1), font_size=12, color=CINZA_ESCURO)

        self._add_textbox(slide, Inches(0.8), Inches(4.2), Inches(5.8), Inches(2.4), "Dor 2: Qualidade", font_size=14, bold=True, color=AZUL_ESCURO)
        self._add_textbox(slide, Inches(0.8), Inches(4.6), Inches(5.8), Inches(1.4), _lines(dor2), font_size=12, color=CINZA_ESCURO)

        self._add_textbox(slide, Inches(7.0), Inches(2.1), Inches(5.8), Inches(2.4), "Dor 3: Produtividade", font_size=14, bold=True, color=AZUL_ESCURO)
        self._add_textbox(slide, Inches(7.0), Inches(2.5), Inches(5.8), Inches(1.8), _lines(dor3), font_size=12, color=CINZA_ESCURO)

        self._add_textbox(slide, Inches(7.0), Inches(4.2), Inches(5.8), Inches(2.4), "Dor 4 e 5", font_size=14, bold=True, color=AZUL_ESCURO)
        self._add_textbox(slide, Inches(7.0), Inches(4.6), Inches(5.8), Inches(0.9), _lines(dor4), font_size=12, color=CINZA_ESCURO)
        self._add_textbox(slide, Inches(7.0), Inches(5.5), Inches(5.8), Inches(1.7), _lines(dor5), font_size=12, color=CINZA_ESCURO)

    def _slide_07_cenario_critico(self, resultados: ResultadosFinanceiros):
        slide = self._add_slide()
        self._add_title_bar(slide, "Cenário Crítico — Visão Geral dos Custos")

        # Grande métrica central
        self._add_metric_box(
            slide, Inches(3.5), Inches(1.8), Inches(6), Inches(1.5),
            "CUSTO TOTAL ANUAL IDENTIFICADO",
            self._fmt(resultados.custo_total_anual_inacao),
            color=VERMELHO,
        )

        # 5 dores
        box_w = Inches(2.5)
        start_x = Inches(0.7)
        gap = Inches(0.3)
        top = Inches(4)

        categorias = [
            ("Dor 1 - Mão de Obra", resultados.total_dor1),
            ("Dor 2 - Qualidade", resultados.total_dor2),
            ("Dor 3 - Produtividade", resultados.total_dor3),
            ("Dor 4 - Segurança", resultados.total_dor4),
            ("Dor 5 - Custos Ocultos", resultados.total_dor5),
        ]

        for i, (label, valor) in enumerate(categorias):
            self._add_metric_box(
                slide, start_x + i * (box_w + gap), top, box_w, Inches(1.2),
                label, self._fmt(valor), color=AZUL_MEDIO,
            )

    def _slide_custos_categoria(self, titulo, breakdown, total):
        """Slide genérico de custos por categoria."""
        slide = self._add_slide()
        self._add_title_bar(slide, titulo)

        # Tabela de breakdown
        rows_data = [["Subcategoria", "Custo Anual (R$)"]]
        for nome, valor in breakdown.items():
            rows_data.append([nome, self._fmt(valor)])
        rows_data.append(["TOTAL", self._fmt(total)])

        n_rows = len(rows_data)
        table = self._add_table(
            slide, Inches(2), Inches(1.8), Inches(9), Inches(0.5 * n_rows + 0.3),
            n_rows, 2, rows_data,
            col_widths=[Inches(5.5), Inches(3.5)],
        )

        # Destacar linha total
        last_row = n_rows - 1
        for c in range(2):
            cell = table.cell(last_row, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = AZUL_ESCURO
            for p in cell.text_frame.paragraphs:
                p.font.bold = True
                p.font.color.rgb = BRANCO

    def _slide_08_custos_operacionais(self, resultados: ResultadosFinanceiros):
        self._slide_custos_categoria(
            "Quantificação — Dor 1: Mão de Obra",
            resultados.breakdown_dor1, resultados.total_dor1,
        )

    def _slide_09_custos_qualidade(self, resultados: ResultadosFinanceiros):
        self._slide_custos_categoria(
            "Quantificação — Dor 2: Qualidade",
            resultados.breakdown_dor2, resultados.total_dor2,
        )

    def _slide_10_custos_seguranca(self, resultados: ResultadosFinanceiros):
        self._slide_custos_categoria(
            "Quantificação — Dor 3: Produtividade",
            resultados.breakdown_dor3, resultados.total_dor3,
        )

    def _slide_11_custos_produtividade(self, resultados: ResultadosFinanceiros):
        self._slide_custos_categoria(
            "Quantificação — Dor 4: Segurança e Ergonomia",
            resultados.breakdown_dor4, resultados.total_dor4,
        )

    def _slide_12_custos_ocultos(self, resultados: ResultadosFinanceiros):
        self._slide_custos_categoria(
            "Quantificação — Dor 5: Custos Ocultos",
            resultados.breakdown_dor5,
            resultados.total_dor5,
        )

    def _slide_13_consolidacao(self, resultados: ResultadosFinanceiros):
        slide = self._add_slide()
        self._add_title_bar(slide, "Consolidação Financeira")

        # Tabela consolidada
        table_data = [
            ["Categoria", "Custo Anual (R$)", "% do Total"],
            [
                "Dor 1 - Mão de Obra",
                self._fmt(resultados.total_dor1),
                self._fmt_pct(
                    resultados.total_dor1 / resultados.custo_total_anual_inacao * 100
                    if resultados.custo_total_anual_inacao > 0 else 0
                ),
            ],
            [
                "Dor 2 - Qualidade",
                self._fmt(resultados.total_dor2),
                self._fmt_pct(
                    resultados.total_dor2 / resultados.custo_total_anual_inacao * 100
                    if resultados.custo_total_anual_inacao > 0 else 0
                ),
            ],
            [
                "Dor 3 - Produtividade",
                self._fmt(resultados.total_dor3),
                self._fmt_pct(
                    resultados.total_dor3 / resultados.custo_total_anual_inacao * 100
                    if resultados.custo_total_anual_inacao > 0 else 0
                ),
            ],
            [
                "Dor 4 - Segurança",
                self._fmt(resultados.total_dor4),
                self._fmt_pct(
                    resultados.total_dor4 / resultados.custo_total_anual_inacao * 100
                    if resultados.custo_total_anual_inacao > 0 else 0
                ),
            ],
            [
                "Dor 5 - Custos Ocultos",
                self._fmt(resultados.total_dor5),
                self._fmt_pct(
                    resultados.total_dor5 / resultados.custo_total_anual_inacao * 100
                    if resultados.custo_total_anual_inacao > 0 else 0
                ),
            ],
            ["TOTAL", self._fmt(resultados.custo_total_anual_inacao), "100.0%"],
        ]

        table = self._add_table(
            slide, Inches(1.5), Inches(1.8), Inches(10), Inches(3.2),
            len(table_data), 3, table_data,
            col_widths=[Inches(4.5), Inches(3), Inches(2.5)],
        )

        # Destacar última linha
        last_row = len(table_data) - 1
        for c in range(3):
            cell = table.cell(last_row, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = AZUL_ESCURO
            for p in cell.text_frame.paragraphs:
                p.font.bold = True
                p.font.color.rgb = BRANCO

        # Métrica de ganho potencial
        self._add_metric_box(
            slide, Inches(3.5), Inches(5.5), Inches(6), Inches(1.2),
            "GANHO ANUAL POTENCIAL COM AUTOMAÇÃO",
            self._fmt(resultados.ganho_anual_potencial),
            color=VERDE,
        )

    def _slide_13_escopo_tecnico(self):
        slide = self._add_slide()
        self._add_title_bar(slide, "Escopo Técnico da Solução")
        self._add_subtitle(slide, "Detalhamento técnico a ser definido em fase de projeto")

        itens = [
            "• Definição de tecnologias e equipamentos",
            "• Layout da célula de automação",
            "• Integração com sistemas existentes (MES/ERP)",
            "• Especificação de sensores e atuadores",
            "• Requisitos de infraestrutura",
            "• Cronograma de implantação",
            "• Plano de comissionamento e start-up",
        ]

        for i, item in enumerate(itens):
            self._add_textbox(
                slide, Inches(1.5), Inches(2.2 + i * 0.55),
                Inches(10), Inches(0.5),
                item, font_size=16, color=CINZA_ESCURO,
            )

    def _slide_14_investimento(self, investimento: InvestimentoAutomacao):
        slide = self._add_slide()
        self._add_title_bar(slide, "Investimento")

        box_w = Inches(3.5)
        gap = Inches(0.5)
        start_x = Inches(1.5)
        top = Inches(2.5)

        self._add_metric_box(
            slide, start_x, top, box_w, Inches(1.3),
            "Investimento Mínimo",
            self._fmt(investimento.valor_investimento_min),
            color=AZUL_MEDIO,
        )
        self._add_metric_box(
            slide, start_x + box_w + gap, top, box_w, Inches(1.3),
            "Investimento Máximo",
            self._fmt(investimento.valor_investimento_max),
            color=AZUL_MEDIO,
        )
        self._add_metric_box(
            slide, start_x + 2 * (box_w + gap), top, box_w, Inches(1.3),
            "Investimento Médio",
            self._fmt(investimento.valor_investimento_medio),
            color=AZUL_ESCURO,
        )

        self._add_textbox(
            slide, Inches(1), Inches(4.8), Inches(11), Inches(1),
            "Nota: Os valores apresentados são estimativas iniciais e podem variar "
            "conforme detalhamento técnico do projeto.",
            font_size=13, color=CINZA_MEDIO, alignment=PP_ALIGN.CENTER,
        )

    def _slide_15_viabilidade(self, resultados: ResultadosFinanceiros,
                              investimento: InvestimentoAutomacao):
        slide = self._add_slide()
        self._add_title_bar(slide, "Viabilidade Financeira")

        # Payback
        payback_txt = (
            f"{resultados.payback_anos:.1f} anos"
            if resultados.payback_anos != float("inf")
            else "N/A"
        )

        box_w = Inches(2.7)
        gap = Inches(0.35)
        start_x = Inches(0.7)
        top = Inches(2)

        self._add_metric_box(
            slide, start_x, top, box_w, Inches(1.3),
            "Payback Simples", payback_txt, color=AZUL_ESCURO,
        )
        self._add_metric_box(
            slide, start_x + (box_w + gap), top, box_w, Inches(1.3),
            "ROI 1 Ano", self._fmt_pct(resultados.roi_1_ano), color=AZUL_MEDIO,
        )
        self._add_metric_box(
            slide, start_x + 2 * (box_w + gap), top, box_w, Inches(1.3),
            "ROI 3 Anos", self._fmt_pct(resultados.roi_3_anos), color=VERDE,
        )
        self._add_metric_box(
            slide, start_x + 3 * (box_w + gap), top, box_w, Inches(1.3),
            "ROI 5 Anos", self._fmt_pct(resultados.roi_5_anos), color=VERDE,
        )

        # Tabela comparativa
        table_data = [
            ["", "Investimento Médio", "Ganho Anual", "Ganho Acumulado"],
            ["Ano 1",
             self._fmt(investimento.valor_investimento_medio),
             self._fmt(resultados.ganho_anual_potencial),
             self._fmt(resultados.ganho_anual_potencial)],
            ["Ano 3", "—",
             self._fmt(resultados.ganho_anual_potencial),
             self._fmt(resultados.ganho_anual_potencial * 3)],
            ["Ano 5", "—",
             self._fmt(resultados.ganho_anual_potencial),
             self._fmt(resultados.ganho_anual_potencial * 5)],
        ]

        self._add_table(
            slide, Inches(1.5), Inches(4), Inches(10), Inches(2),
            4, 4, table_data,
            col_widths=[Inches(1.5), Inches(3), Inches(2.5), Inches(3)],
        )

    def _slide_16_proximas_etapas(self):
        slide = self._add_slide()
        self._add_title_bar(slide, "Próximas Etapas")

        etapas = [
            ("1.", "Validação dos dados e premissas apresentados"),
            ("2.", "Detalhamento técnico do escopo da solução"),
            ("3.", "Proposta comercial formal"),
            ("4.", "Aprovação e kick-off do projeto"),
            ("5.", "Desenvolvimento e implantação"),
            ("6.", "Comissionamento e start-up"),
        ]

        for i, (num, texto) in enumerate(etapas):
            # Número em destaque
            self._add_textbox(
                slide, Inches(2), Inches(2 + i * 0.7),
                Inches(0.5), Inches(0.5),
                num, font_size=20, bold=True, color=AZUL_ESCURO,
            )
            # Texto
            self._add_textbox(
                slide, Inches(2.7), Inches(2 + i * 0.7),
                Inches(8), Inches(0.5),
                texto, font_size=18, color=CINZA_ESCURO,
            )

        # Contato
        self._add_textbox(
            slide, Inches(1), Inches(6.3), Inches(11), Inches(0.5),
            "Estamos à disposição para esclarecer quaisquer dúvidas.",
            font_size=14, color=CINZA_MEDIO, alignment=PP_ALIGN.CENTER,
        )
